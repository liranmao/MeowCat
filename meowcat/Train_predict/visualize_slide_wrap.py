#!/usr/bin/env python3
"""
Generate a PowerPoint with all result images per sample.

Directory layout expected (from the earlier pipeline):
<OUT_ROOT>/<SAMPLE>/
  ├── clusters/
  │     └── <sample>_kmeans_k*.png
  ├── cluster_highlights/          (optional)
  │     └── *.png
  ├── celltype_intensity_percentiles/
  │     └── masked_*.png
  ├── <sample>_predicted_celltype_map.png
  └── summary.txt                  (optional)

Examples:
  python make_slides.py \
    --out-root /project/KidneyHE/output_images \
    --pptx /project/KidneyHE/output_images/results_summary.pptx \
    --title "KidneyHE – Lung Cohort Results"

You can limit to specific samples:
  python make_slides.py --out-root ... --pptx ... --samples P8_LUAD,P12_LUAD
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import argparse
import os
from pathlib import Path
from typing import List, Tuple, Iterable
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ----------------------------------------
# Helpers
# ----------------------------------------
def find_first(patterns: List[str]) -> Path | None:
    for pat in patterns:
        hits = sorted(Path().glob(pat))
        if hits:
            return hits[0]
    return None

def list_pngs(folder: Path) -> List[Path]:
    if not folder.exists():
        return []
    return sorted([p for p in folder.glob("*.png") if p.is_file()])

def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)

def slide_title(slide, text: str, slide_width, top_in=0.1, size_pt=28):
    """
    Add a title textbox across the top of the slide, using the given slide_width
    (pass prs.slide_width from the Presentation).
    """
    left = Inches(0.5)
    top = Inches(top_in)
    width = slide_width - Inches(1)   # 0.5" left margin + 0.5" right margin
    height = Inches(0.6)

    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    p.alignment = PP_ALIGN.LEFT
    return tx


def add_caption(slide, text: str, left, top, width, height=Inches(0.3), size_pt=12, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    p.alignment = align
    return tb

def image_size(path: Path) -> Tuple[int, int]:
    with Image.open(path) as im:
        return im.width, im.height

def layout_grid_positions(slide_w, slide_h, cols, rows, margin_in=0.5, gutter_in=0.2, title_reserved_in=0.9, caption_in=0.35) -> List[Tuple[int,int,int,int,int]]:
    """
    Return positions (left, top, width, height_image, caption_height) in EMU for a grid.
    Reserve space at top for a title.
    """
    margin = Inches(margin_in)
    gutter = Inches(gutter_in)
    reserved = Inches(title_reserved_in)

    avail_w = slide_w - 2*margin - (cols-1)*gutter
    avail_h = slide_h - reserved - 2*margin - (rows-1)*gutter

    cell_w = int(avail_w / cols)
    cell_h = int(avail_h / rows)

    positions = []
    for r in range(rows):
        for c in range(cols):
            left = margin + c*(cell_w + gutter)
            top = reserved + margin + r*(cell_h + gutter)
            # allocate some space at the bottom of the cell for caption
            cap_h = Inches(caption_in)
            img_h = cell_h - cap_h
            positions.append((int(left), int(top), int(cell_w), int(img_h), int(cap_h)))
    return positions

def add_image_fitting(slide, img_path: Path, left, top, box_w, box_h):
    """
    Fit image into a box (preserve aspect), centered vertically.
    Returns the picture shape and its (left, top, width, height).
    """
    iw, ih = image_size(img_path)
    img_aspect = iw / ih
    box_aspect = box_w / box_h

    if img_aspect >= box_aspect:
        # constrained by width
        target_w = box_w
        target_h = int(target_w / img_aspect)
    else:
        # constrained by height
        target_h = box_h
        target_w = int(target_h * img_aspect)

    # center vertically in the box
    dy = (box_h - target_h) // 2
    pic = slide.shapes.add_picture(str(img_path), left, top + dy, width=target_w, height=target_h)
    return pic, (left, top + dy, target_w, target_h)

def batched(iterable: Iterable, n: int) -> List[List]:
    batch = []
    for item in iterable:
        batch.append(item)
        if len(batch) == n:
            yield batch
            batch = []
    if batch:
        yield batch

# ----------------------------------------
# Deck building
# ----------------------------------------
def add_title_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title + subtitle
    slide.shapes.title.text = title
    if subtitle is not None:
        slide.placeholders[1].text = subtitle

def add_overview_slide(prs, sample, cluster_img, argmax_img, summary_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_title(slide, f"{sample} – Overview", prs.slide_width, top_in=0.1, size_pt=28)


    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = Inches(0.5)
    gutter = Inches(0.3)
    title_reserved = Inches(0.9)

    # two columns
    box_w = (slide_w - 2*margin - gutter) // 2
    box_h = slide_h - title_reserved - 2*margin

    left1 = margin
    left2 = margin + box_w + gutter
    top = title_reserved + margin

    if cluster_img and cluster_img.exists():
        add_image_fitting(slide, cluster_img, left1, top, box_w, box_h - Inches(0.35))
        add_caption(slide, "KMeans clusters", left1, top + box_h - Inches(0.35), box_w)

    if argmax_img and argmax_img.exists():
        add_image_fitting(slide, argmax_img, left2, top, box_w, box_h - Inches(0.35))
        add_caption(slide, "Predicted cell-type map (argmax)", left2, top + box_h - Inches(0.35), box_w)

    # Optional summary text box
    if summary_text:
        tb_w = slide_w - 2*margin
        tb_h = Inches(0.6)
        tb_left = margin
        tb_top = slide_h - margin - tb_h
        tb = slide.shapes.add_textbox(tb_left, tb_top, tb_w, tb_h)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = summary_text.strip()
        run.font.size = Pt(12)

def add_gallery(prs: Presentation, sample: str, images: List[Path], title: str, cols=3, rows=2):
    if not images:
        return
    per_slide = cols * rows
    for page, chunk in enumerate(batched(images, per_slide), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        page_suffix = f" – page {page}" if len(images) > per_slide else ""
        slide_title(slide, f"{sample} – {title}{page_suffix}", prs.slide_width, top_in=0.1, size_pt=26)
        
        locs = layout_grid_positions(prs.slide_width, prs.slide_height, cols, rows, title_reserved_in=0.9)
        for (img, (left, top, w, h, cap_h)) in zip(chunk, locs):
            add_image_fitting(slide, img, left, top, w, h)
            add_caption(slide, img.stem, left, top + h, w, cap_h)

# ----------------------------------------
# Main
# ----------------------------------------
def build_deck(out_root: Path,
               pptx_path: Path,
               title: str | None = None,
               samples: List[str] | None = None,
               include_highlights: bool = True,
               intensity_cols: int = 3,
               intensity_rows: int = 2,
               highlight_cols: int = 4,
               highlight_rows: int = 3):
    prs = Presentation()
    if title:
        add_title_slide(prs, title, subtitle=str(out_root))

    # Determine sample list
    candidates = [p.name for p in out_root.iterdir() if p.is_dir()]
    # Prefer those starting with P, but allow explicit --samples override
    if samples:
        sample_list = [s for s in samples if (out_root / s).is_dir()]
    else:
        sample_list = sorted([s for s in candidates if s.startswith("S") or s.startswith("P")])

    for sample in sample_list:
        sdir = out_root / sample
        clusters_dir = sdir / "clusters"
        highlights_dir = sdir / "cluster_highlights"
        intensity_dir = sdir / "celltype_intensity_percentiles"

        # Overview images
        cluster_map = None
        # search for any kmeans cluster map
        cm_hits = list_pngs(clusters_dir)
        if cm_hits:
            cluster_map = cm_hits[0]

        argmax_path = sdir / f"{sample}_predicted_celltype_map.png"
        argmax_map = argmax_path if argmax_path.exists() else None

        summary_text = None
        summary_file = sdir / "summary.txt"
        if summary_file.exists():
            summary_text = summary_file.read_text(encoding="utf-8")

        add_overview_slide(prs, sample, cluster_map, argmax_map, summary_text)

        # Intensity gallery
        intensity_imgs = list_pngs(intensity_dir)
        add_gallery(prs, sample, intensity_imgs, title="Cell-type intensity (masked, percentile-scaled)",
                    cols=intensity_cols, rows=intensity_rows)

        # Highlights gallery (optional)
        if include_highlights:
            highlight_imgs = list_pngs(highlights_dir)
            add_gallery(prs, sample, highlight_imgs, title="Cluster highlights",
                        cols=highlight_cols, rows=highlight_rows)

    ensure_parent(pptx_path)
    prs.save(str(pptx_path))
    print(f"Saved deck: {pptx_path}")

def parse_args():
    ap = argparse.ArgumentParser(description="Create a PowerPoint with result images per sample.")
    ap.add_argument("--out-root", required=True, type=Path,
                    help="Root folder containing per-sample result folders (from the pipeline).")
    ap.add_argument("--pptx", required=True, type=Path,
                    help="Output .pptx path.")
    ap.add_argument("--title", default=None, help="Deck title (optional).")
    ap.add_argument("--samples", default=None,
                    help="Comma-separated list of samples to include (default: all starting with 'P').")
    ap.add_argument("--no-highlights", action="store_true",
                    help="Do not include cluster highlight slides.")
    ap.add_argument("--intensity-cols", type=int, default=3)
    ap.add_argument("--intensity-rows", type=int, default=2)
    ap.add_argument("--highlight-cols", type=int, default=4)
    ap.add_argument("--highlight-rows", type=int, default=3)
    return ap.parse_args()

def main():
    args = parse_args()
    sample_list = None
    if args.samples:
        sample_list = [s.strip() for s in args.samples.split(",") if s.strip()]
    build_deck(
        out_root=args.out_root,
        pptx_path=args.pptx,
        title=args.title,
        samples=sample_list,
        include_highlights=not args.no_highlights,
        intensity_cols=args.intensity_cols,
        intensity_rows=args.intensity_rows,
        highlight_cols=args.highlight_cols,
        highlight_rows=args.highlight_rows,
    )

if __name__ == "__main__":
    main()
